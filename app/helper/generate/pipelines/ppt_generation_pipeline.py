"""PowerPoint generation pipeline for temp-doc service."""

import base64
from io import BytesIO
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ...schemas.temp_doc_schema import ExtractedData, ExtractedPptData


class PptGenerationPipeline:
    """Generate PPTX documents from extracted payloads.

    Highest-fidelity path:
    - If extracted payload includes package XML parts + binary parts/media,
      rebuild the PPTX package directly to preserve themes/background/masters.
    Fallback path:
    - Reconstruct slides from parsed slide content.
    """

    def run(
        self,
        extracted_data: ExtractedData | ExtractedPptData,
        title: str | None = None,
    ) -> bytes:
        """Generate PPTX bytes from structured extracted payload."""
        if isinstance(extracted_data, ExtractedPptData):
            rebuilt = self._try_rebuild_from_package_dump(extracted_data)
            if rebuilt is not None:
                return rebuilt

        prs = Presentation()

        if isinstance(extracted_data, ExtractedPptData):
            self._from_ppt_extracted(prs, extracted_data, title)
        else:
            self._from_json(prs, extracted_data, title)

        output = BytesIO()
        prs.save(output)
        return output.getvalue()

    def _try_rebuild_from_package_dump(self, data: ExtractedPptData) -> bytes | None:
        xml_parts = self._collect_xml_parts(data)
        if not xml_parts:
            return None
        if not self._required_parts_present(xml_parts):
            return None

        media_bytes_by_path = self._collect_media_bytes(data)
        self._collect_slide_media_bytes(data, media_bytes_by_path)
        binary_bytes_by_path = self._collect_binary_bytes(data)

        return self._write_rebuilt_archive(
            xml_parts,
            binary_bytes_by_path,
            media_bytes_by_path,
        )

    def _collect_xml_parts(self, data: ExtractedPptData) -> dict[str, str]:
        """Collect XML package parts keyed by normalized path."""
        parts = getattr(data, "parts", []) or []
        xml_parts: dict[str, str] = {}
        for part in parts:
            if isinstance(part, dict):
                path = (part.get("path") or "").lstrip("/")
                xml = part.get("xml") or ""
            else:
                path = (getattr(part, "path", "") or "").lstrip("/")
                xml = getattr(part, "xml", "") or ""
            if path:
                xml_parts[path] = xml
        return xml_parts

    def _required_parts_present(self, xml_parts: dict[str, str]) -> bool:
        """Validate the minimal XML parts needed to rebuild a pptx package."""
        required = ["[Content_Types].xml",
                    "_rels/.rels", "ppt/presentation.xml"]
        return all(path in xml_parts for path in required)

    def _collect_media_bytes(self, data: ExtractedPptData) -> dict[str, bytes]:
        """Collect media bytes from top-level media entries."""
        media_bytes_by_path: dict[str, bytes] = {}
        for media in (getattr(data, "media", []) or []):
            b64 = self._get(media, "base64_data") or self._get(media, "base64")
            path = (self._get(media, "local_file_path") or "").lstrip("/")
            if not b64 or not path:
                continue
            try:
                media_bytes_by_path[path] = base64.b64decode(b64)
            except (ValueError, TypeError):
                continue
        return media_bytes_by_path

    def _collect_slide_media_bytes(
        self,
        data: ExtractedPptData,
        media_bytes_by_path: dict[str, bytes],
    ) -> None:
        """Collect media bytes from parsed slide picture shapes."""
        for slide in (getattr(data, "parsed_slides", []) or []):
            for shape in (self._get(slide, "shapes", []) or []):
                self._maybe_add_shape_media_bytes(shape, media_bytes_by_path)

    def _maybe_add_shape_media_bytes(
        self,
        shape,
        media_bytes_by_path: dict[str, bytes],
    ) -> None:
        """Add picture shape base64 bytes to media map when valid and missing."""
        if self._get(shape, "kind") != "picture":
            return
        path = (self._get(shape, "target_path") or "").lstrip("/")
        b64 = self._get(shape, "base64")
        if not path or not b64 or path in media_bytes_by_path:
            return
        try:
            media_bytes_by_path[path] = base64.b64decode(b64)
        except (ValueError, TypeError):
            return

    def _collect_binary_bytes(self, data: ExtractedPptData) -> dict[str, bytes]:
        """Collect non-XML binary package parts from base64 payloads."""
        binary_bytes_by_path: dict[str, bytes] = {}
        for item in (getattr(data, "binary_parts", []) or []):
            path = (self._get(item, "path") or "").lstrip("/")
            b64 = self._get(item, "base64")
            if not path or not b64:
                continue
            try:
                binary_bytes_by_path[path] = base64.b64decode(b64)
            except (ValueError, TypeError):
                continue
        return binary_bytes_by_path

    def _write_rebuilt_archive(
        self,
        xml_parts: dict[str, str],
        binary_bytes_by_path: dict[str, bytes],
        media_bytes_by_path: dict[str, bytes],
    ) -> bytes:
        """Write a reconstructed pptx zip archive from XML/binary/media parts."""
        output = BytesIO()
        with ZipFile(output, "w", compression=ZIP_DEFLATED) as archive:
            for path, xml in xml_parts.items():
                archive.writestr(path, xml)

            for path, blob in binary_bytes_by_path.items():
                if path in media_bytes_by_path:
                    continue
                archive.writestr(path, blob)

            for path, blob in media_bytes_by_path.items():
                archive.writestr(path, blob)

        return output.getvalue()

    def _from_ppt_extracted(self, prs: Presentation, data: ExtractedPptData, title: str | None) -> None:
        if title:
            self._add_title_slide(prs, title, "")

        slides = getattr(data, "slides", []) or []
        paragraphs = getattr(data, "paragraphs", []) or []
        tables = getattr(data, "tables", []) or []
        media = getattr(data, "media", []) or []
        parsed_slides = getattr(data, "parsed_slides", []) or []

        if slides:
            paragraphs_by_idx = self._index_items_by_field(paragraphs, "index")
            tables_by_idx = self._index_items_by_field(tables, "index")
            media_by_idx = dict(enumerate(media))

            self._render_structured_slides(
                prs,
                slides,
                paragraphs_by_idx,
                tables_by_idx,
                media_by_idx,
            )
            return

        if parsed_slides:
            self._render_parsed_slides(prs, parsed_slides)
            return

        if paragraphs or tables:
            self._from_json(
                prs,
                ExtractedData(
                    document_order=getattr(data, "document_order", []) or [],
                    styles=getattr(data, "styles", []) or [],
                    paragraphs=paragraphs,
                    tables=tables,
                    media=media,
                ),
                title,
            )

    def _index_items_by_field(self, items: list, field: str) -> dict:
        """Index dict/object items by a field resolved via _get()."""
        indexed: dict = {}
        for item in items:
            idx = self._get(item, field)
            if idx is not None:
                indexed[idx] = item
        return indexed

    def _slide_sort_key(self, slide) -> int:
        """Stable sort key for slide-like payload items."""
        return self._get(slide, "index", 10**9)

    def _resolve_structured_slide_title(self, slide) -> str:
        """Resolve title for structured slide model."""
        fallback_num = self._get(slide, "slide_number")
        if fallback_num is None:
            fallback_num = self._get(slide, "index", 0) + 1
        return (self._get(slide, "title") or f"Slide {fallback_num}").strip()

    def _resolve_parsed_slide_title(self, slide) -> str:
        """Resolve title for parsed slide model."""
        fallback_num = self._get(slide, "index", 0) + 1
        return (self._get(slide, "title") or f"Slide {fallback_num}").strip()

    def _structured_slide_lines(self, slide, paragraphs_by_idx: dict, slide_title: str) -> list[str]:
        """Build cleaned body lines for a structured slide."""
        lines: list[str] = []
        for p_idx in (self._get(slide, "paragraph_indices", []) or []):
            para = paragraphs_by_idx.get(p_idx)
            if para is None:
                continue
            text = (self._get(para, "text") or "").strip()
            if text:
                lines.append(text)

        if not lines:
            fallback_text = (self._get(slide, "text") or "").strip()
            if fallback_text:
                lines = [line.strip()
                         for line in fallback_text.splitlines() if line.strip()]

        return self._clean_body_lines(lines, slide_title)

    def _tables_from_indices(self, slide, tables_by_idx: dict) -> list[list[list[str]]]:
        """Build table matrices from slide table index references."""
        slide_tables: list[list[list[str]]] = []
        for t_idx in (self._get(slide, "table_indices", []) or []):
            table = tables_by_idx.get(t_idx)
            if table is None:
                continue
            rows = self._rows_from_generic_table(table)
            if rows:
                slide_tables.append(rows)
        return slide_tables

    def _rows_from_generic_table(self, table) -> list[list[str]]:
        """Extract row/cell text matrix from table-like object/dict."""
        rows: list[list[str]] = []
        for row in (self._get(table, "rows", []) or []):
            rows.append(
                [
                    (self._get(cell, "text") or "")
                    for cell in (self._get(row, "cells", []) or [])
                ]
            )
        return rows

    def _structured_slide_media(self, slide, media_by_idx: dict) -> list:
        """Collect non-placeholder media entries referenced by slide."""
        slide_media = []
        for m_idx in (self._get(slide, "media_indices", []) or []):
            media = media_by_idx.get(m_idx)
            if media is not None and not self._is_placeholder_media(media):
                slide_media.append(media)
        return slide_media

    def _render_structured_slides(
        self,
        prs: Presentation,
        slides: list,
        paragraphs_by_idx: dict,
        tables_by_idx: dict,
        media_by_idx: dict,
    ) -> None:
        """Render slides from structured ppt extraction model."""
        for slide in sorted(slides, key=self._slide_sort_key):
            slide_title = self._resolve_structured_slide_title(slide)
            slide_lines = self._structured_slide_lines(
                slide,
                paragraphs_by_idx,
                slide_title,
            )
            slide_tables = self._tables_from_indices(slide, tables_by_idx)
            slide_media = self._structured_slide_media(slide, media_by_idx)
            title_color_rgb = self._pick_title_color_from_paragraphs(
                slide_title,
                self._get(slide, "paragraph_indices", []) or [],
                paragraphs_by_idx,
            )
            out_slide = self._add_composite_slide(
                prs,
                title=slide_title or "Slide",
                lines=slide_lines,
                tables=slide_tables,
                media_items=slide_media,
                title_color_rgb=title_color_rgb,
            )
            self._pad_to_shape_count(
                out_slide, self._get(slide, "shape_count"))

    def _parsed_slide_lines(self, slide, slide_title: str) -> list[str]:
        """Build cleaned body lines for parsed slide model, including notes."""
        slide_lines: list[str] = []
        slide_text = (self._get(slide, "text") or "").strip()
        if slide_text:
            slide_lines = [line.strip()
                           for line in slide_text.splitlines() if line.strip()]
            slide_lines = self._clean_body_lines(slide_lines, slide_title)

        notes = self._get(slide, "notes")
        if isinstance(notes, dict):
            notes_text = (notes.get("text") or "").strip()
            if notes_text:
                slide_lines.append(f"Notes: {notes_text}")
        return slide_lines

    def _parsed_slide_tables(self, slide) -> list[list[list[str]]]:
        """Build table matrices from parsed slide shapes."""
        slide_tables: list[list[list[str]]] = []
        for shape in (self._get(slide, "shapes", []) or []):
            if self._get(shape, "kind") != "graphic_frame":
                continue
            if self._get(shape, "graphic_type") != "table":
                continue
            table_payload = self._get(shape, "table") or {}
            rows = self._rows_from_generic_table(table_payload)
            if rows:
                slide_tables.append(rows)
        return slide_tables

    def _render_parsed_slides(self, prs: Presentation, parsed_slides: list) -> None:
        """Render slides from parsed_slides fallback model."""
        for slide in sorted(parsed_slides, key=self._slide_sort_key):
            if self._get(slide, "parse_error"):
                continue
            slide_title = self._resolve_parsed_slide_title(slide)
            slide_lines = self._parsed_slide_lines(slide, slide_title)
            slide_tables = self._parsed_slide_tables(slide)
            out_slide = self._add_composite_slide(
                prs,
                title=slide_title or "Slide",
                lines=slide_lines,
                tables=slide_tables,
                media_items=[],
                title_color_rgb=self._pick_title_color_from_shapes(
                    self._get(slide, "shapes", []) or [],
                ),
            )
            self._pad_to_shape_count(
                out_slide, self._get(slide, "shape_count"))

    def _from_json(self, prs: Presentation, data: ExtractedData, title: str | None) -> None:
        if title:
            self._add_title_slide(prs, title, "")

        para_by_idx = {p.index: p for p in data.paragraphs}
        table_by_idx = {t.index: t for t in data.tables}

        if data.document_order:
            self._render_json_in_order(prs, data, para_by_idx, table_by_idx)
            return

        self._render_json_sorted(prs, data)

    def _render_json_in_order(
        self,
        prs: Presentation,
        data: ExtractedData,
        para_by_idx: dict,
        table_by_idx: dict,
    ) -> None:
        """Render slides from generic extracted JSON respecting document_order."""
        for item in data.document_order:
            if item.type == "paragraph":
                paragraph = para_by_idx.get(item.index)
                self._render_json_paragraph(prs, paragraph)
            elif item.type == "table":
                table = table_by_idx.get(item.index)
                self._render_json_table(prs, table)

    def _render_json_sorted(self, prs: Presentation, data: ExtractedData) -> None:
        """Render slides from generic extracted JSON using index order."""
        for paragraph in sorted(data.paragraphs, key=lambda item: item.index):
            self._render_json_paragraph(prs, paragraph)
        for table in sorted(data.tables, key=lambda item: item.index):
            self._render_json_table(prs, table)

    def _render_json_paragraph(self, prs: Presentation, paragraph) -> None:
        """Render a single text paragraph as a slide when non-empty."""
        if paragraph is None:
            return
        text = paragraph.text or ""
        if not text.strip():
            return
        self._add_text_slide(prs, title=(
            paragraph.style or "Paragraph"), lines=[text])

    def _render_json_table(self, prs: Presentation, table) -> None:
        """Render a single table as a table slide."""
        if table is None:
            return
        rows = [[(cell.text or "") for cell in row.cells]
                for row in table.rows]
        self._add_table_slide(prs, title="Table", rows=rows)

    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str) -> None:
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

    def _add_text_slide(self, prs: Presentation, title: str, lines: list[str]) -> None:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title

        body = slide.shapes.placeholders[1].text_frame
        body.clear()

        first = True
        for line in lines:
            text = (line or "").strip()
            if not text:
                continue

            if first:
                p = body.paragraphs[0]
                first = False
            else:
                p = body.add_paragraph()

            p.text = text
            p.level = 0
            p.alignment = PP_ALIGN.LEFT
            if p.runs:
                p.runs[0].font.size = Pt(18)

    def _add_table_slide(self, prs: Presentation, title: str, rows: list[list[str]]) -> None:
        layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title

        if not rows:
            return

        col_count = max((len(r) for r in rows), default=0)
        if col_count == 0:
            return

        norm_rows = [r + [""] * (col_count - len(r)) for r in rows]

        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9.0)
        height = Inches(5.0)

        table_shape = slide.shapes.add_table(
            rows=len(norm_rows),
            cols=col_count,
            left=left,
            top=top,
            width=width,
            height=height,
        )
        table = table_shape.table

        for r_i, row in enumerate(norm_rows):
            for c_i, value in enumerate(row):
                table.cell(r_i, c_i).text = value or ""

    def _add_composite_slide(
        self,
        prs: Presentation,
        title: str,
        lines: list[str],
        tables: list[list[list[str]]],
        media_items: list,
        title_color_rgb: str | None = None,
    ):
        layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title
            self._apply_title_color(slide.shapes.title, title_color_rgb)

        text_lines = [ln for ln in lines if (ln or "").strip()]

        if text_lines:
            self._add_composite_textbox(
                slide,
                text_lines,
                has_media=bool(media_items),
                has_table=bool(tables),
            )

        if media_items:
            self._add_composite_media(slide, media_items)

        if tables:
            self._add_table_shape(
                slide,
                rows=tables[0],
                left=Inches(0.5),
                top=Inches(3.9),
                width=Inches(9.0),
                height=Inches(2.8),
            )

        return slide

    def _add_composite_textbox(
        self,
        slide,
        text_lines: list[str],
        *,
        has_media: bool,
        has_table: bool,
    ) -> None:
        """Add the text area to a composite slide."""
        left = Inches(0.5)
        top = Inches(1.2)
        text_width = Inches(6.2 if has_media else 9.0)
        text_height = Inches(2.4 if has_table else 5.4)

        textbox = slide.shapes.add_textbox(left, top, text_width, text_height)
        tf = textbox.text_frame
        tf.clear()

        first = True
        for line in text_lines:
            paragraph = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            paragraph.text = (line or "").strip()
            paragraph.level = 0
            paragraph.alignment = PP_ALIGN.LEFT
            if paragraph.runs:
                paragraph.runs[0].font.size = Pt(16)

    def _add_composite_media(self, slide, media_items: list) -> None:
        """Add right-rail media stack to a composite slide."""
        media_left = Inches(7.0)
        media_top = Inches(1.2)
        max_w = Inches(2.4)
        max_h = Inches(1.6)

        for media in media_items[:3]:
            added = self._add_media_image(
                slide, media, media_left, media_top, max_w, max_h)
            if added:
                media_top += Inches(1.8)

    def _apply_title_color(self, title_shape, color_rgb: str | None) -> None:
        color = self._normalize_hex_rgb(color_rgb)
        if not color:
            return
        try:
            tf = title_shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = RGBColor.from_string(color)
        except (ValueError, AttributeError, TypeError):
            return

    def _normalize_hex_rgb(self, value: str | None) -> str | None:
        if not value:
            return None
        v = value.strip()
        if v.lower().startswith("scheme:"):
            return None
        if v.startswith("#"):
            v = v[1:]
        if len(v) != 6:
            return None
        try:
            int(v, 16)
        except ValueError:
            return None
        return v.upper()

    def _pick_title_color_from_paragraphs(
        self,
        slide_title: str,
        paragraph_indices: list[int],
        paragraphs_by_idx: dict[int, object],
    ) -> str | None:
        """Pick title color from title-matching paragraph runs, then any runs."""
        color = self._color_from_indexed_paragraphs(
            paragraph_indices,
            paragraphs_by_idx,
            title_filter=(slide_title or "").strip().lower(),
        )
        if color:
            return color
        return self._color_from_indexed_paragraphs(
            paragraph_indices,
            paragraphs_by_idx,
            title_filter=None,
        )

    def _color_from_indexed_paragraphs(
        self,
        paragraph_indices: list[int],
        paragraphs_by_idx: dict[int, object],
        title_filter: str | None,
    ) -> str | None:
        """Return first valid run color from indexed paragraphs."""
        for idx in paragraph_indices:
            para = paragraphs_by_idx.get(idx)
            if para is None:
                continue
            if not self._paragraph_matches_title_filter(para, title_filter):
                continue
            color = self._first_valid_run_color(
                self._get(para, "runs", []) or [])
            if color:
                return color
        return None

    def _paragraph_matches_title_filter(self, para, title_filter: str | None) -> bool:
        """Check whether paragraph text matches title filter if provided."""
        if not title_filter:
            return True
        text = (self._get(para, "text") or "").strip()
        return bool(text and text.lower() == title_filter)

    def _first_valid_run_color(self, runs: list) -> str | None:
        """Return the first valid normalized color from run list."""
        for run in runs:
            color = self._get(run, "color_rgb")
            if self._normalize_hex_rgb(color):
                return color
        return None

    def _pick_title_color_from_shapes(self, shapes: list[dict]) -> str | None:
        for shape in shapes or []:
            if not shape.get("is_title"):
                continue
            for para in shape.get("paragraphs", []) or []:
                for run in para.get("runs", []) or []:
                    color = run.get("color_rgb")
                    if self._normalize_hex_rgb(color):
                        return color
        return None

    def _pad_to_shape_count(self, slide, expected_count: int | None) -> None:
        if expected_count is None:
            return
        missing = expected_count - len(slide.shapes)
        if missing <= 0:
            return

        for _ in range(missing):
            shp = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(9.8),
                Inches(7.2),
                Inches(0.01),
                Inches(0.01),
            )
            shp.fill.background()
            shp.line.fill.background()

    def _add_table_shape(self, slide, rows: list[list[str]], left, top, width, height) -> None:
        if not rows:
            return
        col_count = max((len(r) for r in rows), default=0)
        if col_count == 0:
            return
        norm_rows = [r + [""] * (col_count - len(r)) for r in rows]
        table_shape = slide.shapes.add_table(
            rows=len(norm_rows),
            cols=col_count,
            left=left,
            top=top,
            width=width,
            height=height,
        )
        table = table_shape.table
        for r_i, row in enumerate(norm_rows):
            for c_i, value in enumerate(row):
                table.cell(r_i, c_i).text = value or ""

    def _add_media_image(self, slide, media, left, top, max_w, max_h) -> bool:
        img_blob = None
        b64 = self._get(media, "base64_data") or self._get(media, "base64")
        if b64:
            try:
                img_blob = base64.b64decode(b64)
            except (ValueError, TypeError):
                img_blob = None

        if not img_blob:
            return False

        try:
            stream = BytesIO(img_blob)
            slide.shapes.add_picture(
                stream, left, top, width=max_w, height=max_h)
            return True
        except (ValueError, TypeError):
            return False

    def _clean_body_lines(self, lines: list[str], title: str) -> list[str]:
        out: list[str] = []
        title_norm = (title or "").strip().lower()
        seen: set[str] = set()
        for ln in lines:
            t = (ln or "").strip()
            if not t:
                continue
            t_norm = t.lower()
            if title_norm and t_norm == title_norm:
                continue
            if t_norm in seen:
                continue
            seen.add(t_norm)
            out.append(t)
        return out

    def _is_placeholder_media(self, media) -> bool:
        src = self._get(media, "source")
        name = ""
        if isinstance(src, dict):
            name = (src.get("name") or "").lower()
        if "placeholder" in name:
            return True
        return False

    def _get(self, obj, key: str, default=None):
        if isinstance(obj, dict):
            return obj.get(key, default)
        return getattr(obj, key, default)
